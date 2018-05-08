#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
PPTX 文件读取示例

https://python-pptx.readthedocs.io/en/latest/index.html#api-documentation

python-pptx has the following capabilities, with many more on the roadmap:

Round-trip any Open XML presentation (.pptx file) including all its elements
Add slides
Populate text placeholders, for example to create a bullet slide
Add image to slide at arbitrary position and size
Add textbox to a slide; manipulate text font size and bold
Add table to a slide
Add auto shapes (e.g. polygons, flowchart shapes, etc.) to a slide
Add and manipulate column, bar, line, and pie charts
Access and change core document properties such as title and subject
"""

import os
import pptx
from pptx.util import Inches

from common import loadcfgcm
from common import filecm
from common import imagecm
from common import logcm
from common import opencvcm

# 配置
default_config = """
{
    "img_root": "./images/faces/",
    "max_count" : 10,
    "fps": 1,
    "save_path": "./temp/output.pptx",
    "width": 1024,
    "height": 683
}
"""

# 加载配置文件
cfg = loadcfgcm.load("file_pptx_from_img.json", default_config)

pptFile = pptx.Presentation()

# 取得图片列表
path_list = filecm.search_files(cfg["img_root"], '.jpeg', r'^[^\.]+')

# 按图片编号顺序导入
for fn in path_list:
    slide = pptFile.slides.add_slide(pptFile.slide_layouts[1])

    # 为PPTX文件当前幻灯片中第一个文本框设置文字，本文代码中可忽略
    slide.shapes.placeholders[0].text = fn[:fn.rindex('.')]

    # 导入并为当前幻灯片添加图片，起始位置和尺寸可修改
    slide.shapes.add_picture(fn, Inches(0), Inches(0), Inches(10), Inches(7.5))

pptFile.save('./temp/img_to_pptx.pptx')
